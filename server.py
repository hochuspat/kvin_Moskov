from fastapi import FastAPI, File, UploadFile, HTTPException
from pydantic import BaseModel
from PIL import Image
import pytesseract
import io
import nltk
import spacy
from pptx import Presentation
import random
import whisper
import gensim
from fastapi.responses import FileResponse
import pytextrank
from fastapi.middleware.cors import CORSMiddleware
from typing import List
from docx import Document
from docx.shared import Pt
import os
app = FastAPI()
model = whisper.load_model("large-v3")

@app.post("/upload-audio")
async def upload_audio(file: UploadFile = File(...)):
    try:
        # Сохранение аудиофайла
        with open("temp_audio.mp3", "wb") as buffer:
            buffer.write(file.file.read())

        # Использование Whisper для распознавания текста
        result = model.transcribe("temp_audio.mp3")

        # Возвращаем распознанный текст
        return {"text": result["text"]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

nlp = spacy.load("ru_core_news_sm")
nlp.add_pipe("textrank")

class Text(BaseModel):
    text: str

@app.post("/upload-image")
async def create_upload_file(file: UploadFile = File(...)):
    contents = await file.read()
    image = Image.open(io.BytesIO(contents))
    text = pytesseract.image_to_string(image, lang='rus')
    return {"text": text}
def clean_text(text):
    return "".join(ch for ch in text if ch.isprintable() or ch in '\n\r\t')

def create_presentation(text):
    prs = Presentation()
    # Здесь можно попытаться выбрать случайный стиль, но функционал ограничен
    # Например, выбор случайного макета слайда (если доступно несколько макетов)
    slide_layout = random.choice(prs.slide_layouts)

    for paragraph in text.split('\n\n'):  # Разделение текста на абзацы
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Лекция"  # Заголовок слайда
        content = slide.placeholders[1]  # Предполагаем, что это место для основного контента
        content.text = paragraph

    prs.save('lecture.pptx')
@app.post("/process-text")
async def process_text(text: Text):
    try:
        sentences = nltk.sent_tokenize(text.text)

        # Обработка коротких текстов
        if len(sentences) < 3:
            return {"summary": text.text}

        doc = nlp(text.text)
        keywords = [p.text for p in doc._.phrases[:10]]
        dictionary = gensim.corpora.Dictionary([keywords])
        corpus = [dictionary.doc2bow(nltk.word_tokenize(s)) for s in sentences]
        lda = gensim.models.LdaModel(corpus, num_topics=5, id2word=dictionary)
        topics = [lda.get_document_topics(bow) for bow in corpus]
        max_topics = [max(t, key=lambda x: x[1])[0] for t in topics]

        groups = {}
        for i, t in enumerate(max_topics):
            if t not in groups:
                groups[t] = []
            if len(groups[t]) < 5: 
                groups[t].append(sentences[i])

        summary = ""
        for t in sorted(groups):
            summary += f"## Глава {t+1}: {', '.join(keywords[:3])}\n\n"
            summary += " ".join(groups[t]) + "\n\n"
        number_of_chapters = len(groups)
        # Очистка текста
        clean_ocr_text = clean_text(text.text)
        clean_summary = clean_text(summary)

        # Создание файла Word
        doc = Document()
        doc.add_heading('Лекция', level=1)
        doc.add_paragraph(clean_ocr_text)
        doc.add_heading('Конспект', level=1)
        doc.add_paragraph(clean_summary)

        # Сохранение файла
        file_path = "lecture.docx"
        doc.save(file_path)
        create_presentation(clean_summary)

        # Возвращаем путь к файлу
        return {"summary": summary,"chapters": number_of_chapters, "chapters": len(groups), "wordFilePath": file_path}
    except Exception as e:
        print(e)
        raise HTTPException(status_code=500, detail=str(e))
# Добавьте эту секцию для настройки CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Разрешает все источники
    allow_credentials=True,
    allow_methods=["*"],  # Разрешает все методы
    allow_headers=["*"],  # Разрешает все заголовки
)

class ChapterPerformance(BaseModel):
    chapter: str
    performance: int

class LectureData(BaseModel):
    summary: str
    presentation: str
    wordFile: str
    audioFile: str
    statistics: List[ChapterPerformance]

# Заглушка данных для примера
example_lecture_data = LectureData(
    summary='Конспект лекции...',
    presentation='Ссылка на презентацию',
    wordFile='Ссылка на Word файл',
    audioFile='Ссылка на аудио файл',
    statistics=[
        ChapterPerformance(chapter='Глава 1', performance=75),
        ChapterPerformance(chapter='Глава 2', performance=60),
        # Дополнительные главы...
    ],
)

@app.get("/lecture-data", response_model=LectureData)
async def get_lecture_data():
    return example_lecture_data
# Пример структуры данных проходимости
performance_data = [
    {"chapter": "Глава 1", "performance": 75},
    {"chapter": "Глава 2", "performance": 60},
    # Дополнительные главы...
]


@app.get("/performance")
async def get_performance():
    return performance_data

@app.get("/download-word/{file_path}")
async def download_word(file_path: str):
    # Проверка безопасности пути к файлу

    return FileResponse(file_path, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', filename="lecture.docx")
@app.get("/download-presentation")
async def download_presentation():
    file_path = "lecture.pptx"  # Путь к файлу презентации
    return FileResponse(file_path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename="lecture.pptx")
